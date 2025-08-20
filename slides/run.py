from pptx import Presentation

# Carrega o modelo existente
prs = Presentation("modelo.pptx")

# Usa o primeiro layout do modelo
layout = prs.slide_layouts[0]

# Cria um novo slide baseado no layout
slide = prs.slides.add_slide(layout)

# Preenche título e subtítulo
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Minha Apresentação Automática"
subtitle.text = "Gerada com Python e modelo PPTX"

# Salva como novo arquivo
prs.save("apresentacao_final.pptx")
